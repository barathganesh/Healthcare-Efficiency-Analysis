"""
Step 20 — Fix & Polish the Final Presentation
================================================
Addresses Gemini's review + data-accuracy fixes:

CRITICAL FIXES (Judge Traps):
  A. Nursing Paradox  — reframe from "too few nurses" to utilization/mix
  B. Admin Contradiction — remove as root cause (p=0.170)
  C. Occupancy Buffer  — clarify phrasing (surge capacity not occupancy level)

DATA ACCURACY FIXES:
  D. Slide 10 cost breakdown — align with actual Fix 1-5 categories
  E. Slide 17 state numbers — replace fabricated TX/FL with actual top states

NARRATIVE POLISH:
  F. Slide titles & hooks per Gemini's consulting-style suggestions
  G. Conclusion & closing script
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import os, copy

# ── Paths ──────────────────────────────────────────────────────────
SRC = '/Users/barathganesh/Downloads/Healthcare_Analytics_FINAL_PRESENTATION.pptx'
OUT = '/Users/barathganesh/Documents/ResearchProjects/healthcareAnalytics/outputs/Healthcare_Analytics_FINAL_v2.pptx'

prs = Presentation(SRC)
slides = list(prs.slides)

# ── Helper: replace text in a slide, preserving run-level formatting ─
def replace_text(slide, old, new, *, verbose=True):
    """Replace `old` with `new` inside any shape on `slide`.
    Tries run-level first; falls back to paragraph-level."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            # --- Try within individual runs first ---
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)
                    if verbose:
                        print(f'  ✓ run-level: "{old[:50]}…" → replaced')
                    return True
            # --- Fall back: paragraph-level (concat all runs) ---
            full = ''.join(r.text for r in para.runs)
            if old in full:
                new_full = full.replace(old, new)
                if para.runs:
                    para.runs[0].text = new_full
                    for r in para.runs[1:]:
                        r.text = ''
                if verbose:
                    print(f'  ✓ para-level: "{old[:50]}…" → replaced')
                return True
    if verbose:
        print(f'  ✗ NOT FOUND: "{old[:60]}…"')
    return False


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — The Crisis → The Burning Platform
# ══════════════════════════════════════════════════════════════════
print('\n── Slide 2: Crisis → Burning Platform ──')
replace_text(slides[1],
    'The Crisis: Hospitals at Risk',
    'The Burning Platform: Operational Inefficiency is Costing Lives')
replace_text(slides[1],
    'hospitals struggling',
    'hospitals at risk')


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — Finding 1: Size
# ══════════════════════════════════════════════════════════════════
print('\n── Slide 6: Size Is the Enemy → Diseconomies of Scale ──')
replace_text(slides[5],
    'Finding 1: Size Is the Enemy',
    'Finding 1: Diseconomies of Scale — Bigger Isn\'t Better')
replace_text(slides[5],
    'IMPLICATION: STOP consolidating hospitals',
    'IMPLICATION: Stop incentivizing consolidation — smaller is more efficient')


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — Finding 2: Occupancy Paradox
# ══════════════════════════════════════════════════════════════════
print('\n── Slide 7: Occupancy — clarify buffer ──')
replace_text(slides[6],
    'IMPLICATION: Maintain 35-40% buffer capacity',
    'IMPLICATION: Maintain 60%+ surge capacity (hero-level occupancy ≈35-40%)')


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — Finding 3: Admin — already correct ("Myth Busted")
# ══════════════════════════════════════════════════════════════════
print('\n── Slide 8: Admin Myth — strengthening language ──')
replace_text(slides[7],
    'IMPLICATION: Focus on size/occupancy, not admin',
    'IMPLICATION: Admin is NOT the problem — focus on size, occupancy & revenue')


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — Finding 4: Revenue
# ══════════════════════════════════════════════════════════════════
print('\n── Slide 9: Revenue pivot ──')
replace_text(slides[8],
    'ACTION: Optimize revenue cycle, not just cut costs',
    'Strategic Pivot: Revenue Integrity Over Cost Cutting')


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — Financial Impact: fix breakdown to match actual data
# ══════════════════════════════════════════════════════════════════
print('\n── Slide 10: Financial breakdown ──')
# Actual fix categories from total_impact_summary.csv:
#   Fix 5 Readmission Savings:  $1.43B (58%)
#   Fix 3 Occupancy Revenue:    $0.81B (33%)
#   Fix 1 Admin Reduction:      $0.20B (8%)
#   Fix 4 Contract Labor:       $0.06B (2%)
#   Fix 2 Nursing Investment:  -$0.06B (cost)
replace_text(slides[9],
    'Labor: $1.2B (49%)',
    'Readmission Savings: $1.43B (58%)')
replace_text(slides[9],
    'Supply chain: $0.6B (25%)',
    'Revenue Optimization: $0.81B (33%)')
replace_text(slides[9],
    'Overhead: $0.4B (16%)',
    'Admin Efficiency: $0.20B (8%)')
replace_text(slides[9],
    'Revenue cycle: $0.24B (10%)',
    'Contract Labor + Nursing Net: $0.00B (1%)')


# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — Clinical Impact (numbers are approximately correct)
# ══════════════════════════════════════════════════════════════════
# No changes needed — 76,182 total is correct, condition split sums correctly


# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — Root Causes: FIX THE TWO JUDGE TRAPS
# ══════════════════════════════════════════════════════════════════
print('\n── Slide 12: Root Causes — fix admin + nursing contradictions ──')

# FIX A: Remove "too many admin, too few nurses"
#   → Admin p=0.170 (not significant per slide 8)
#   → Nursing: Heroes have MORE hrs/bed (279 vs 218), so 109 hospitals need
#     staffing MIX optimization, not blanket "more nurses"
replace_text(slides[11],
    '4. WORKFORCE MISALLOCATION - Too many admin, too few nurses',
    '4. WORKFORCE OPTIMIZATION - Staffing mix misaligned vs hero benchmarks (109 hospitals)')

# Also clarify #2 to be more precise
replace_text(slides[11],
    '2. CAPACITY MISMANAGEMENT - Wrong occupancy levels',
    '2. CAPACITY MISMANAGEMENT - Congestion from wrong occupancy levels')


# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — Short-Term Recs: fix occupancy phrasing
# ══════════════════════════════════════════════════════════════════
print('\n── Slide 13: Short-term — fix occupancy phrasing ──')
replace_text(slides[12],
    'reduce occupancy to 35-40%',
    'maintain 60%+ surge capacity')
# Strengthen non-safety-net wording
replace_text(slides[12],
    'Rebalance workforce, increase volume',
    'Optimize staffing mix, increase volume')


# ══════════════════════════════════════════════════════════════════
# SLIDE 14 — Long-Term Recs
# ══════════════════════════════════════════════════════════════════
print('\n── Slide 14: Long-term recs ──')
replace_text(slides[13],
    'Right-size to 60-80 beds, build buffer capacity',
    'Right-size to 60-80 beds, build 60%+ surge capacity')
replace_text(slides[13],
    'Stop incentivizing consolidation',
    'End consolidation incentives')
replace_text(slides[13],
    'Combine efficiency + financial support',
    'Pair efficiency mandates with financial support')


# ══════════════════════════════════════════════════════════════════
# SLIDE 15 — Implementation Plan
# ══════════════════════════════════════════════════════════════════
print('\n── Slide 15: Implementation ──')
replace_text(slides[14],
    'PHASE 3 (Months 13-36): Sustainability & Scale',
    'PHASE 3 (Months 13-36): Sustainability & Market Leadership')
replace_text(slides[14],
    'Deepen, expand, transform, lead',
    'Deepen analytics, expand best practices, build sustainability')


# ══════════════════════════════════════════════════════════════════
# SLIDE 17 — Geographic Impact: fix fabricated state numbers
# ══════════════════════════════════════════════════════════════════
print('\n── Slide 17: Geographic — fix state numbers ──')
# Actual top 5 by lives saved:
#   CA: 8,054 | NY: 6,560 | AR: 5,491 | GA: 5,314 | MO: 3,808
# Actual top 5 by savings:
#   GA: $194M | CA: $192M | NY: $163M | VA: $141M | AR: $136M
replace_text(slides[16],
    'California: 8,200 | Texas: 7,800',
    'California: 8,054 | New York: 6,560')
replace_text(slides[16],
    'Florida: 6,400 | New York: 5,900',
    'Arkansas: 5,491 | Georgia: 5,314')
replace_text(slides[16],
    'California: $312M | Texas: $298M',
    'Georgia: $194M | California: $192M')
replace_text(slides[16],
    'Florida: $245M | New York: $226M',
    'New York: $163M | Virginia: $141M')


# ══════════════════════════════════════════════════════════════════
# SLIDE 18 — Why This Analysis Wins
# ══════════════════════════════════════════════════════════════════
print('\n── Slide 18: Why This Wins ──')
replace_text(slides[17],
    '3. MYTH-BUSTING: Size, occupancy, admin findings',
    '3. MYTH-BUSTING: Size ≠ strength, occupancy ≠ revenue, admin ≠ root cause')


# ══════════════════════════════════════════════════════════════════
# SLIDE 19 — Conclusion: sharper hook
# ══════════════════════════════════════════════════════════════════
print('\n── Slide 19: Conclusion hook ──')
replace_text(slides[18],
    "Not 'Can we?' but 'Will we?'",
    "Not 'Can we afford to act?' but 'Can we afford not to?'")
replace_text(slides[18],
    "Our answer is YES. What's yours?",
    'The model works. The path is clear.')


# ══════════════════════════════════════════════════════════════════
# SLIDE 20 — Thank You: stronger close
# ══════════════════════════════════════════════════════════════════
print('\n── Slide 20: Closing script ──')
replace_text(slides[19],
    'The data is clear.',
    '163,000 hospital-years of evidence.')
replace_text(slides[19],
    'The path is mapped.',
    '232 ready-to-implement prescriptions.')
replace_text(slides[19],
    'The tools are ready.',
    '$2.44 billion in savings, quantified.')
replace_text(slides[19],
    "Let's save 76,182 lives.",
    "76,182 lives are waiting. Let's get to work.")


# ══════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f'\n{"="*60}')
print(f'✅  Saved: {OUT}')
print(f'{"="*60}')
